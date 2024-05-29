import multiprocessing
import os
import glob
import shutil
import threading
import time
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from enum import Enum
from concurrent.futures import ThreadPoolExecutor, as_completed

output_temp_dir_default = '../temp/pdf2ppt.output/JPEG'


class Result(Enum):
    UnknownError = -1,
    Success = 0,
    FileNotExistError = 1,
    InvalidIndexError = 2,


def scan_pdf_file() -> tuple[Result, list[str]]:
    path_pattern: str = os.path.join("../", "*.pdf")
    pdf_path_list: list[str] = glob.glob(path_pattern)
    pdf_filename_list: list[str] = [os.path.basename(path) for path in pdf_path_list]

    pdf_file_num = len(pdf_path_list)

    if pdf_file_num == 0:
        console_output("Error: No PDF file.")
        return Result.FileNotExistError, pdf_path_list

    if pdf_file_num == 1:
        return Result.Success, pdf_path_list

    if pdf_file_num > 1:
        console_output("Multiple PDF files exist.")
        index = 0
        for pdf_filename in pdf_filename_list:
            console_output(f"[{index}] {pdf_filename}")
            index += 1
        user_input = input("Process all these files? (y/n) ")
        user_input_lowercase = user_input.lower()
        if user_input_lowercase == "y" or user_input_lowercase == "yes":
            return Result.Success, pdf_path_list
        else:
            console_output("Choose one PDF file by index.")
            console_output('Example: To choose "[5] example.pdf", you should input "5".')
            input_index: str = input("Index: ")
            input_index_num: int = int(input_index)
            if input_index_num < 0 or input_index_num > pdf_file_num - 1:
                console_output("Invalid index.")
                return Result.InvalidIndexError, pdf_path_list
            return Result.Success, [pdf_path_list[input_index_num]]


def convert_pdf_to_ppt_multi(pdf_path_list: list[str], output_temp_dir: str, jpeg_quality: int = 100) -> Result:
    if len(pdf_path_list) == 0:
        return Result.FileNotExistError

    file_str: str = "file" if len(pdf_path_list) == 1 else "files"
    console_output(f"Task PDF to PPT start. Total {len(pdf_path_list)} {file_str}.")

    if len(pdf_path_list) == 1:
        return convert_pdf_to_ppt(pdf_path_list[0], output_temp_dir, jpeg_quality)

    results: list[Result] = []
    create_dir_if_not_exists(output_temp_dir)

    with ThreadPoolExecutor(max_workers=len(pdf_path_list)) as executor:
        future_to_pdf = {executor.submit(convert_pdf_to_ppt, pdf, output_temp_dir, jpeg_quality): pdf for pdf in
                         pdf_path_list}

        for future in as_completed(future_to_pdf):
            pdf = future_to_pdf[future]
            result = future.result()
            results.append(result)

    if all(isinstance(result, Result) and result == Result.Success for result in results):
        return Result.Success
    else:
        return Result.UnknownError


def convert_pdf_to_ppt(pdf_path: str, output_temp_dir_base: str, jpeg_quality: int = 100) -> Result:
    doc = fitz.open(pdf_path)
    doc_len = len(doc)

    filename: str = os.path.basename(pdf_path)
    output_temp_dir: str = f"{output_temp_dir_base}/{filename}"
    create_dir_if_not_exists(output_temp_dir)

    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    output_pdf_file_to_img(doc, pdf_path, filename, output_temp_dir, jpeg_quality)

    console_output(f"[{filename}] Process complete. Merging start.")
    for i in range(doc_len):
        img_output_path = get_img_path(output_temp_dir, i + 1)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 使用空白页模板

        slide.shapes.add_picture(img_output_path, 0, 0, width=presentation.slide_width,
                                 height=presentation.slide_height)
    console_output(f"[{filename}] Merging complete. Saving...")

    pptx_path = f"{pdf_path}.pptx"
    presentation.save(pptx_path)

    console_output(f"Complete! Save to: {pptx_path}.")
    return Result.Success


class ProcessCounter:
    def __init__(self, section_name: str, total: int, log_tick: int = 2) -> None:
        self.__section_name: str = section_name
        self.__total: int = total
        self.__current: int = 0
        self.__log_tick: int = log_tick
        self.__last_log: int = 0

    def add(self, count: int) -> None:
        self.__current += count
        if self.__current >= self.__total:
            self.log()
        else:
            if self.__current - self.__last_log >= self.__log_tick:
                self.log()

    def log(self):
        self.__last_log = self.__current
        percentage_str: str = f"{(self.__current / self.__total) * 100:3.0f}".strip().rjust(3)
        console_output(
            f"[{self.__section_name}] [{percentage_str}%] Complete {self.__current}/{self.__total}.")


def output_pdf_file_to_img(doc: fitz.Document, filepath: str, filename: str, output_temp_dir: str, jpeg_quality):
    doc_len: int = len(doc)
    num_processes: int = 8
    pages_per_process: int = int(doc_len / num_processes)

    counter: ProcessCounter = ProcessCounter(filename, doc_len, log_tick=1)
    counter.log()

    events = [multiprocessing.Event() for _ in range(num_processes)]
    queue = multiprocessing.Queue()

    processes = []
    for index in range(num_processes):
        start_index = index * pages_per_process
        end_index = (index + 1) * pages_per_process - 1 if index != num_processes - 1 else doc_len - 1
        p = multiprocessing.Process(
            target=output_pdf_range_to_img,
            args=(filepath, output_temp_dir, jpeg_quality, start_index, end_index, events[index], queue)
        )
        processes.append(p)
        p.start()

    monitor_thread = threading.Thread(target=monitor_progress, args=(queue, counter))
    monitor_thread.start()

    for p in processes:
        p.join()

    for event in events:
        event.wait()

    queue.put(None)
    monitor_thread.join()


def monitor_progress(queue: multiprocessing.Queue, counter: ProcessCounter) -> None:
    while True:
        message = queue.get()
        if message is None:
            break
        counter.add(message)


def output_pdf_range_to_img(filepath: str, output_temp_dir: str, jpeg_quality, index_start: int,
                            index_end: int, thread_event: multiprocessing.Event, queue: multiprocessing.Queue) -> None:
    doc = fitz.open(filepath)
    doc_len = len(doc)

    def process_one_page(index: int) -> None:
        output_pdf_page_as_jpeg(doc, index, get_img_path(output_temp_dir, index + 1), jpeg_quality)
        queue.put(1)  # Notify the main process
        thread_event.set()

    with ThreadPoolExecutor(max_workers=int(os.cpu_count() / 2)) as executor:
        futures = [executor.submit(process_one_page, i) for i in range(index_start, index_end + 1)]
        img_paths = [future.result() for future in as_completed(futures)]


def get_img_path(dir_base: str, page_number: int) -> str:
    return f"{dir_base}/page_{page_number}.jpg"


def output_pdf_page_as_jpeg(doc: fitz.Document, index: int, output_file_jpeg: str, jpeg_quality: int) -> None:
    page = doc.load_page(index)

    zoom = 3840 / page.rect.width
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)

    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

    img.save(output_file_jpeg, "JPEG", quality=jpeg_quality)


def create_dir_if_not_exists(dir_path: str) -> None:
    if not os.path.exists(dir_path):
        os.makedirs(dir_path)


def app_main() -> None:
    file_scan_result, pdf_path_list = scan_pdf_file()
    if file_scan_result != Result.Success:
        return

    app_start_time = time.time()
    output_temp_dir_base: str = output_temp_dir_default

    console_output("Preparing cache...")
    if os.path.exists(output_temp_dir_base):
        shutil.rmtree(output_temp_dir_base)
    os.makedirs(output_temp_dir_base)
    console_output("Prepare cache complete.")

    convert_result = convert_pdf_to_ppt_multi(pdf_path_list, output_temp_dir_default, jpeg_quality=100)
    app_end_time = time.time()
    execution_time = app_end_time - app_start_time

    if convert_result == Result.Success:
        console_output(f"All task complete. Using {execution_time:.2f}s.")
    else:
        console_output(f"Not all tasks success. Using {execution_time:.2f}s.")


console_output_lock = threading.Lock()


def console_output(message: str) -> None:
    with console_output_lock:
        print(message)
